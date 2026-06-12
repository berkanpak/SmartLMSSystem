import os
import fitz
from pptx import Presentation


MAX_CHARS_PER_DOC = 500_000


def extract_pdf_text(file_path: str) -> str:
    doc = fitz.open(file_path)
    try:
        total_chars = 0
        parts = []
        for page in doc:
            text = page.get_text()
            parts.append(text)
            total_chars += len(text)
            if total_chars > MAX_CHARS_PER_DOC:
                parts.append("\n[TRUNCATED: Document exceeds character limit for study assistant]")
                break
        return "".join(parts)
    finally:
        doc.close()


def extract_pptx_text(file_path: str) -> str:
    prs = Presentation(file_path)
    parts = []
    total_chars = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                parts.append(text)
                total_chars += len(text)
        if total_chars > MAX_CHARS_PER_DOC:
            parts.append("\n[TRUNCATED: Presentation exceeds character limit for study assistant]")
            break
    return "\n".join(parts)


def extract_document_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_pdf_text(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx_text(file_path)
    return ""
