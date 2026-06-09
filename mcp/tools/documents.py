import os
import fitz
from pptx import Presentation


def extract_pdf_text(file_path: str) -> str:
    doc = fitz.open(file_path)
    try:
        parts = [page.get_text() for page in doc]
        return "".join(parts)
    finally:
        doc.close()


def extract_pptx_text(file_path: str) -> str:
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def extract_document_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_pdf_text(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx_text(file_path)
    return ""
